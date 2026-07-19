from __future__ import annotations

import argparse
import json
import zipfile
from datetime import date
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
SKILL_DIR = Path(__file__).resolve().parent
DEFAULT_OUTPUT_DIR = ROOT / "outputs" / "growth-menu"


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def check(condition: bool, name: str, details: str, results: list[dict[str, Any]]) -> None:
    results.append({"name": name, "status": "pass" if condition else "fail", "details": details})


def pptx_text_and_fonts(path: Path) -> tuple[str, set[str]]:
    texts: list[str] = []
    fonts: set[str] = set()
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            for elem in root.iter():
                tag = elem.tag.rsplit("}", 1)[-1]
                if tag == "t" and elem.text:
                    texts.append(elem.text)
                if tag in {"latin", "ea", "cs"}:
                    typeface = elem.attrib.get("typeface")
                    if typeface:
                        fonts.add(typeface)
    return "\n".join(texts), fonts


def audit(output_dir: Path, expected_date: str | None = None) -> dict[str, Any]:
    config = load_json(SKILL_DIR / "config.json")
    latest_json = output_dir / "latest.json"
    latest_pptx = output_dir / "latest.pptx"
    notification = output_dir / "latest_notification.md"
    results: list[dict[str, Any]] = []

    check(latest_json.exists(), "latest_json_exists", str(latest_json), results)
    check(latest_pptx.exists(), "latest_pptx_exists", str(latest_pptx), results)
    check(notification.exists(), "notification_exists", str(notification), results)
    if not latest_json.exists() or not latest_pptx.exists():
        return {"status": "fail", "checked_date": expected_date, "checks": results}

    data = load_json(latest_json)
    target_date = expected_date or data.get("date")
    check(data.get("date") == target_date, "date_matches", f"actual={data.get('date')} expected={target_date}", results)

    prs = Presentation(latest_pptx)
    check(
        len(prs.slides) == config["output"]["slide_count"],
        "slide_count",
        f"actual={len(prs.slides)} expected={config['output']['slide_count']}",
        results
    )

    ppt_text, fonts = pptx_text_and_fonts(latest_pptx)
    for forbidden in config["curriculum"]["explicit_exclusions"]:
        check(forbidden.lower() not in ppt_text.lower(), f"excluded_{forbidden}", "not present", results)
    for forbidden in config["output"]["forbidden_deck_phrases"]:
        check(forbidden not in ppt_text, f"not_submission_form_{forbidden}", "not present in PowerPoint", results)
    check(config["output"]["font_family"] in fonts, "font_family", f"fonts={sorted(fonts)}", results)

    required_teaching_markers = ["基本概念", "ポイント1", "ケースの状況", "理論で読み解くと", "今日の型", "ルール"]
    for marker in required_teaching_markers:
        check(marker in ppt_text, f"teaching_marker_{marker}", "present", results)

    check(len(data.get("business_lesson", {}).get("learning_points", [])) == 3,
          "three_learning_points", "exactly 3", results)
    check(len(data.get("business_lesson", {}).get("cause_effect_chain", [])) == 3,
          "three_cause_effect_steps", "exactly 3", results)
    check(len(data.get("case_lesson", {}).get("observations", [])) == 3,
          "three_case_observations", "exactly 3", results)
    check(len(data.get("case_lesson", {}).get("interpretation", [])) == 3,
          "three_case_interpretations", "exactly 3", results)
    check(len(data.get("english_lesson", {}).get("vocabulary", [])) == 5,
          "five_vocabulary_items", "exactly 5", results)

    generation = data.get("generation", {})
    provider = generation.get("provider")
    check(provider in {"offline", "gemini_free_tier"}, "allowed_provider", str(provider), results)
    check("openai" not in json.dumps(generation, ensure_ascii=False).lower(),
          "openai_not_used", "no OpenAI API provider", results)
    check(config["provider"]["allow_paid_fallback"] is False,
          "paid_fallback_disabled", "paid fallback is disabled", results)
    check(config["provider"]["allow_search_grounding"] is False,
          "search_grounding_disabled", "no paid/search grounding", results)

    privacy = config["privacy"]
    feedback = load_json(SKILL_DIR / "feedback" / "latest.json")
    privacy_ok = (
        feedback.get("privacy_check", {}).get("raw_output_stored") is False
        and feedback.get("privacy_check", {}).get("confidential_details_removed") is True
    )
    check(privacy_ok, "feedback_privacy", "raw output is not stored in the public repository", results)
    check(privacy["repository_is_public"] is True, "public_repo_awareness", "privacy guard enabled", results)

    notification_text = notification.read_text(encoding="utf-8") if notification.exists() else ""
    check("latest.pptx" in notification_text, "stable_pptx_link", "notification links to latest.pptx", results)
    check("このチャット" in notification_text, "feedback_return_path", "submission returns to ChatGPT", results)
    check("今日理解したこと" in notification_text, "notification_submission_guide",
          "submission instructions stay outside PowerPoint", results)

    failures = [item for item in results if item["status"] == "fail"]
    return {
        "status": "pass" if not failures else "fail",
        "checked_date": target_date,
        "generated_at": date.today().isoformat(),
        "checks": results,
        "failure_count": len(failures)
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Audit generated HOS morning lesson outputs")
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    parser.add_argument("--date", help="Expected content date in YYYY-MM-DD")
    parser.add_argument("--write-report", action="store_true")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    report = audit(args.output_dir, args.date)
    if args.write_report:
        audit_dir = args.output_dir / "audit"
        audit_dir.mkdir(parents=True, exist_ok=True)
        (audit_dir / "latest.json").write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        if report.get("checked_date"):
            (audit_dir / f"{report['checked_date']}.json").write_text(
                json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
            )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
