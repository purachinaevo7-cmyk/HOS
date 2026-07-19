from __future__ import annotations

import argparse
import json
import os
import re
import sys
from datetime import date, datetime
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from zoneinfo import ZoneInfo

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SKILL_DIR = Path(__file__).resolve().parent
DEFAULT_OUTPUT_DIR = ROOT / "outputs" / "growth-menu"

NAVY = RGBColor(28, 42, 62)
BLUE = RGBColor(58, 87, 122)
PALE = RGBColor(239, 243, 247)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(35, 42, 50)
MUTED = RGBColor(88, 99, 110)
LINE = RGBColor(207, 216, 225)


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def jst_today() -> date:
    return datetime.now(ZoneInfo("Asia/Tokyo")).date()


def read_recent_topics(output_dir: Path, limit: int = 7) -> list[str]:
    history = output_dir / "history"
    if not history.exists():
        return []
    topics: list[str] = []
    for path in sorted(history.glob("*.json"), reverse=True):
        try:
            data = load_json(path)
            theme = str(data.get("business", {}).get("theme", "")).strip()
            if theme:
                topics.append(theme)
        except (OSError, json.JSONDecodeError):
            continue
        if len(topics) >= limit:
            break
    return topics


def offline_fixture(target_date: date, feedback: dict[str, Any]) -> dict[str, Any]:
    next_focus = feedback.get("business", {}).get("next_focus", [])[:2]
    grammar = feedback.get("english", {}).get("next_grammar") or "主語と動詞の骨格"
    return {
        "date": target_date.isoformat(),
        "title": f"{target_date.month}月{target_date.day}日の成長メニュー",
        "focus_summary": "業界構造と個社戦略を分け、根拠から結論までを短くつなぐ",
        "estimated_minutes": 30,
        "previous_feedback_summary": " / ".join(next_focus) if next_focus else "初回のため基礎設計から開始",
        "business": {
            "title": "経営・人事：利益構造と人材配置をつなぐ",
            "theme": "利益構造と人材配置",
            "source_label": "サンプル：企業公式IRページ",
            "source_url": "https://example.com/ir",
            "why_now": "売上成長だけでなく、固定費と人材配置が収益性に与える影響を考えるため。",
            "key_points": [
                "売上成長を数量・単価・構成に分ける",
                "固定費と変動費を分ける",
                "必要な人材能力と配置を財務影響へつなぐ"
            ],
            "question": "利益率を改善するには、顧客価値の向上とコスト構造の改善のどちらを先に打つべきか。",
            "output_task": "結論、根拠2つ、反対意見への留保を150〜250字で書く。"
        },
        "english": {
            "title": f"英語：{grammar}",
            "grammar_pattern": "The most important thing is that S + V.",
            "bad_example": "Most important thing efficient organization.",
            "correct_example": "The most important thing is that the organization allocates talent efficiently.",
            "explanation": "英語は主語と動詞を明示します。thing の後ろは is を置き、内容は that節で続けます。",
            "vocabulary": [
                {"word": "allocate", "meaning": "配分する", "example": "The company allocates capital carefully."},
                {"word": "capability", "meaning": "組織能力", "example": "Digital capability supports growth."},
                {"word": "margin", "meaning": "利益率", "example": "The operating margin improved."},
                {"word": "workforce", "meaning": "労働力・人員", "example": "The workforce needs new skills."},
                {"word": "priority", "meaning": "優先事項", "example": "Productivity is a strategic priority."}
            ],
            "writing_prompt": "『人材配置が持続的成長に重要である理由』を、今日の型を使って英語2文で書く。"
        },
        "notification": {
            "headline": f"{target_date.month}月{target_date.day}日の朝教材ができました",
            "summary": "今日は、利益構造と人材配置をつなげて考えます。英語は主語と動詞の骨格を練習します。",
            "duration": "約30分"
        }
    }


def system_prompt() -> str:
    return (SKILL_DIR / "prompts" / "generator_system.md").read_text(encoding="utf-8")


def user_prompt(target_date: date, config: dict[str, Any], feedback: dict[str, Any], recent_topics: list[str]) -> str:
    shape = {
        "date": "YYYY-MM-DD",
        "title": "M月D日の成長メニュー",
        "focus_summary": "string",
        "estimated_minutes": 20,
        "previous_feedback_summary": "string",
        "business": {
            "title": "string",
            "theme": "string",
            "source_label": "string",
            "source_url": "https://...",
            "why_now": "string",
            "key_points": ["string", "string", "string"],
            "question": "string",
            "output_task": "string"
        },
        "english": {
            "title": "string",
            "grammar_pattern": "string",
            "bad_example": "string",
            "correct_example": "string",
            "explanation": "string",
            "vocabulary": [
                {"word": "string", "meaning": "string", "example": "string"}
            ],
            "writing_prompt": "string"
        },
        "notification": {
            "headline": "string",
            "summary": "string",
            "duration": "約N分"
        }
    }
    return (
        f"対象日: {target_date.isoformat()}\n"
        f"設定: {json.dumps(config, ensure_ascii=False)}\n"
        f"前回評価（公開リポジトリ用に匿名化済み）: {json.dumps(feedback, ensure_ascii=False)}\n"
        f"直近テーマ: {json.dumps(recent_topics, ensure_ascii=False)}\n"
        "次のJSON形状に厳密に合わせてください。vocabularyは必ず5件、key_pointsは必ず3件です。\n"
        f"{json.dumps(shape, ensure_ascii=False, indent=2)}"
    )


def extract_output_text(payload: dict[str, Any]) -> str:
    texts: list[str] = []
    for item in payload.get("output", []):
        if item.get("type") != "message":
            continue
        for content in item.get("content", []):
            if content.get("type") == "output_text" and isinstance(content.get("text"), str):
                texts.append(content["text"])
    if not texts and isinstance(payload.get("output_text"), str):
        texts.append(payload["output_text"])
    if not texts:
        raise RuntimeError("OpenAI response did not contain output text")
    return "\n".join(texts)


def parse_json_text(text: str) -> dict[str, Any]:
    cleaned = text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start == -1 or end == -1 or end <= start:
            raise
        return json.loads(cleaned[start : end + 1])


def call_openai(target_date: date, config: dict[str, Any], feedback: dict[str, Any], recent_topics: list[str]) -> dict[str, Any]:
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is required unless --offline is used")
    model = os.getenv("OPENAI_MODEL", "gpt-5-mini").strip() or "gpt-5-mini"
    base_url = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").rstrip("/")
    body = {
        "model": model,
        "store": False,
        "tools": [{"type": "web_search"}],
        "input": [
            {"role": "system", "content": [{"type": "input_text", "text": system_prompt()}]},
            {"role": "user", "content": [{"type": "input_text", "text": user_prompt(target_date, config, feedback, recent_topics)}]}
        ]
    }
    response = requests.post(
        f"{base_url}/responses",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json=body,
        timeout=180
    )
    if not response.ok:
        raise RuntimeError(f"OpenAI API error {response.status_code}: {response.text[:1000]}")
    return parse_json_text(extract_output_text(response.json()))


def validate_content(data: dict[str, Any], config: dict[str, Any], target_date: date) -> None:
    schema = load_json(SKILL_DIR / "schema.json")
    missing = [key for key in schema["required"] if key not in data]
    if missing:
        raise ValueError(f"Missing root fields: {missing}")
    if data["date"] != target_date.isoformat():
        raise ValueError(f"Date mismatch: {data['date']} != {target_date.isoformat()}")
    for key in schema["business_required"]:
        if key not in data["business"]:
            raise ValueError(f"Missing business field: {key}")
    for key in schema["english_required"]:
        if key not in data["english"]:
            raise ValueError(f"Missing english field: {key}")
    if len(data["business"].get("key_points", [])) != 3:
        raise ValueError("business.key_points must contain exactly 3 items")
    if len(data["english"].get("vocabulary", [])) != 5:
        raise ValueError("english.vocabulary must contain exactly 5 items")
    minutes = int(data["estimated_minutes"])
    out = config["output"]
    if not out["estimated_minutes_min"] <= minutes <= out["estimated_minutes_max"]:
        raise ValueError("estimated_minutes outside configured range")
    full_text = json.dumps(data, ensure_ascii=False).lower()
    for excluded in config["curriculum"]["explicit_exclusions"]:
        if excluded.lower() in full_text:
            raise ValueError(f"Excluded curriculum item found: {excluded}")
    source_url = data["business"]["source_url"]
    parsed = urlparse(source_url)
    if parsed.scheme != "https" or not parsed.netloc:
        raise ValueError("business.source_url must be a valid HTTPS URL")


def set_text(shape, text: str, font_name: str, font_size: int, *, bold: bool = False,
             color: RGBColor = DARK, align: PP_ALIGN = PP_ALIGN.LEFT,
             margin: float = 0.12, valign: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    east_asia = r_pr.find(qn("a:ea"))
    if east_asia is None:
        east_asia = OxmlElement("a:ea")
        r_pr.append(east_asia)
    east_asia.set("typeface", font_name)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_card(slide, x: float, y: float, w: float, h: float, text: str, config: dict[str, Any],
             *, fill: RGBColor = WHITE, line: RGBColor = LINE, font_size: int | None = None,
             bold: bool = False, color: RGBColor = DARK) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    set_text(shape, text, config["output"]["font_family"], font_size or config["output"]["body_pt"], bold=bold, color=color)
    return shape


def add_header(slide, title: str, subtitle: str, config: dict[str, Any]) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    set_text(band, title, config["output"]["font_family"], config["output"]["title_pt"], bold=True, color=WHITE, margin=0.2, valign=MSO_ANCHOR.MIDDLE)
    sub = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(0.12), Inches(4.7), Inches(0.45))
    sub.fill.background()
    sub.line.fill.background()
    set_text(sub, subtitle, config["output"]["font_family"], config["output"]["note_pt"], color=WHITE, align=PP_ALIGN.RIGHT, margin=0.02, valign=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text: str, config: dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(7.08), Inches(12.63), Inches(0.25))
    shape.fill.background()
    shape.line.fill.background()
    set_text(shape, text, config["output"]["font_family"], config["output"]["note_pt"], color=MUTED, margin=0.01, valign=MSO_ANCHOR.MIDDLE)


def render_pptx(data: dict[str, Any], config: dict[str, Any], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_header(slide, data["title"], "HOS Morning Growth Menu", config)
    add_card(slide, 0.45, 1.02, 12.43, 0.88, f"今日の重点\n{data['focus_summary']}", config, fill=PALE, line=BLUE, bold=True, font_size=config["output"]["closing_pt"])
    add_card(slide, 0.45, 2.12, 5.92, 2.1, f"前回から引き継ぐこと\n{data['previous_feedback_summary']}", config)
    add_card(slide, 6.62, 2.12, 6.26, 2.1, "今日の完了条件\n・経営・人事の問いに150〜250字で回答\n・英作文を2〜3文で回答\n・根拠を最低2つ示す", config)
    add_card(slide, 0.45, 4.52, 12.43, 1.43, f"所要時間：約{data['estimated_minutes']}分\n読むより、短く考えて書くことを優先する。", config, fill=WHITE, line=LINE, font_size=config["output"]["closing_pt"], bold=True)
    add_footer(slide, "提出先：このChatGPTスレッド。回答後、評価を匿名化して翌日の教材へ反映。", config)

    business = data["business"]
    slide = prs.slides.add_slide(blank)
    add_header(slide, business["title"], business["theme"], config)
    add_card(slide, 0.45, 0.98, 4.1, 1.35, f"なぜ見るか\n{business['why_now']}", config, fill=PALE, line=BLUE)
    points = "\n".join(f"{i+1}. {p}" for i, p in enumerate(business["key_points"]))
    add_card(slide, 4.8, 0.98, 8.08, 1.35, f"見るポイント\n{points}", config)
    add_card(slide, 0.45, 2.58, 12.43, 1.28, f"考える問い\n{business['question']}", config, fill=WHITE, line=BLUE, bold=True, font_size=config["output"]["closing_pt"])
    add_card(slide, 0.45, 4.12, 12.43, 1.25, f"アウトプット課題\n{business['output_task']}", config, fill=PALE, line=LINE, bold=True)
    source = add_card(slide, 0.45, 5.66, 12.43, 0.82, f"参照：{business['source_label']}\n{business['source_url']}", config, font_size=config["output"]["note_pt"])
    for paragraph in source.text_frame.paragraphs:
        for run in paragraph.runs:
            if business["source_url"] in run.text:
                run.hyperlink.address = business["source_url"]
    add_footer(slide, "公開情報のみを使用。顧客・候補者・社内限定情報は書かない。", config)

    english = data["english"]
    slide = prs.slides.add_slide(blank)
    add_header(slide, english["title"], "Business English", config)
    add_card(slide, 0.45, 0.98, 12.43, 0.75, f"今日の型：{english['grammar_pattern']}", config, fill=PALE, line=BLUE, bold=True, font_size=config["output"]["closing_pt"])
    add_card(slide, 0.45, 1.98, 6.05, 1.42, f"誤り例\n{english['bad_example']}", config)
    add_card(slide, 6.83, 1.98, 6.05, 1.42, f"正解例\n{english['correct_example']}", config, fill=WHITE, line=BLUE)
    add_card(slide, 0.45, 3.65, 12.43, 0.82, f"なぜ：{english['explanation']}", config, font_size=config["output"]["note_pt"])
    vocab_text = " / ".join(f"{v['word']}（{v['meaning']}）" for v in english["vocabulary"])
    add_card(slide, 0.45, 4.72, 12.43, 0.83, f"今日の5語：{vocab_text}", config, fill=PALE, line=LINE, font_size=config["output"]["note_pt"])
    add_card(slide, 0.45, 5.8, 12.43, 0.82, f"英作文：{english['writing_prompt']}", config, fill=WHITE, line=BLUE, bold=True)
    add_footer(slide, "提出は日本語アウトプットと一緒でよい。英語は短くても、主語と動詞を明示する。", config)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def write_outputs(data: dict[str, Any], config: dict[str, Any], output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "history").mkdir(parents=True, exist_ok=True)
    date_slug = data["date"]
    latest_json = output_dir / "latest.json"
    history_json = output_dir / "history" / f"{date_slug}.json"
    latest_pptx = output_dir / "latest.pptx"
    history_pptx = output_dir / "history" / f"{date_slug}_growth_menu.pptx"
    notification = output_dir / "latest_notification.md"

    text = json.dumps(data, ensure_ascii=False, indent=2) + "\n"
    latest_json.write_text(text, encoding="utf-8")
    history_json.write_text(text, encoding="utf-8")
    render_pptx(data, config, latest_pptx)
    history_pptx.write_bytes(latest_pptx.read_bytes())
    raw_url = "https://raw.githubusercontent.com/purachinaevo7-cmyk/HOS/main/outputs/growth-menu/latest.pptx"
    notification.write_text(
        f"## {data['notification']['headline']}\n\n"
        f"{data['notification']['summary']}\n\n"
        f"所要時間：{data['notification']['duration']}\n\n"
        f"[今日のPowerPoint教材を開く]({raw_url})\n\n"
        "学習後は、このチャットに日本語アウトプットと英作文を返してください。\n",
        encoding="utf-8"
    )
    return {
        "latest_json": latest_json,
        "history_json": history_json,
        "latest_pptx": latest_pptx,
        "history_pptx": history_pptx,
        "notification": notification
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate the HOS morning growth-menu PowerPoint")
    parser.add_argument("--date", help="Target date in YYYY-MM-DD; defaults to today in Asia/Tokyo")
    parser.add_argument("--offline", action="store_true", help="Use deterministic fixture and skip OpenAI API")
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    target_date = date.fromisoformat(args.date) if args.date else jst_today()
    config = load_json(SKILL_DIR / "config.json")
    feedback = load_json(SKILL_DIR / "feedback" / "latest.json")
    recent_topics = read_recent_topics(args.output_dir)
    data = offline_fixture(target_date, feedback) if args.offline else call_openai(target_date, config, feedback, recent_topics)
    validate_content(data, config, target_date)
    outputs = write_outputs(data, config, args.output_dir)
    print(json.dumps({"status": "generated", "date": target_date.isoformat(), "outputs": {k: str(v) for k, v in outputs.items()}}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:
        print(f"growth-menu generation failed: {exc}", file=sys.stderr)
        raise
