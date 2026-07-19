from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

NAVY = RGBColor(28, 42, 62)
BLUE = RGBColor(58, 87, 122)
PALE = RGBColor(239, 243, 247)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(35, 42, 50)
MUTED = RGBColor(88, 99, 110)
LINE = RGBColor(207, 216, 225)
ACCENT = RGBColor(104, 126, 153)
GREEN = RGBColor(234, 243, 238)


def set_text(
    shape,
    text: str,
    font_name: str,
    font_size: int,
    *,
    bold: bool = False,
    color: RGBColor = DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    margin: float = 0.12,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    east_asia = r_pr.find(qn("a:ea"))
    if east_asia is None:
        east_asia = OxmlElement("a:ea")
        r_pr.append(east_asia)
    east_asia.set("typeface", font_name)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    config: dict[str, Any],
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    font_size: int | None = None,
    bold: bool = False,
    color: RGBColor = DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    set_text(
        shape,
        text,
        config["output"]["font_family"],
        font_size or config["output"]["body_pt"],
        bold=bold,
        color=color,
        align=align,
    )
    return shape


def add_header(slide, title: str, subtitle: str, config: dict[str, Any]) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    set_text(
        band,
        title,
        config["output"]["font_family"],
        config["output"]["title_pt"],
        bold=True,
        color=WHITE,
        margin=0.2,
        valign=MSO_ANCHOR.MIDDLE,
    )
    sub = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.15),
        Inches(0.12),
        Inches(4.85),
        Inches(0.45),
    )
    sub.fill.background()
    sub.line.fill.background()
    set_text(
        sub,
        subtitle,
        config["output"]["font_family"],
        config["output"]["note_pt"],
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        margin=0.02,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, text: str, config: dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35),
        Inches(7.08),
        Inches(12.63),
        Inches(0.25),
    )
    shape.fill.background()
    shape.line.fill.background()
    set_text(
        shape,
        text,
        config["output"]["font_family"],
        config["output"]["note_pt"],
        color=MUTED,
        margin=0.01,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_arrow(slide, x: float, y: float, w: float = 0.58, h: float = 0.38) -> None:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()


def render_pptx(data: dict[str, Any], config: dict[str, Any], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    profile = data["case_profile"]

    slide = prs.slides.add_slide(blank)
    add_header(slide, data["title"], "HOS Morning Lesson | Context first", config)
    add_card(
        slide,
        0.45,
        1.02,
        12.43,
        0.92,
        f"今日のテーマ\n{data['focus_summary']}",
        config,
        fill=PALE,
        line=BLUE,
        bold=True,
        font_size=config["output"]["closing_pt"],
    )
    lesson = data["business_lesson"]
    add_card(
        slide,
        0.45,
        2.20,
        5.95,
        1.62,
        f"なぜ学ぶか\n{lesson['why_it_matters']}",
        config,
    )
    goals = "\n".join(f"{i + 1}. {item}" for i, item in enumerate(lesson["learning_points"]))
    add_card(
        slide,
        6.67,
        2.20,
        6.21,
        1.62,
        f"今日わかるようになること\n{goals}",
        config,
        fill=WHITE,
        line=BLUE,
    )
    add_card(
        slide,
        0.45,
        4.16,
        12.43,
        1.02,
        f"前回からつなぐ視点\n{data['previous_feedback_summary']}",
        config,
        fill=GREEN,
        line=LINE,
    )
    add_card(
        slide,
        0.45,
        5.48,
        12.43,
        0.80,
        f"会社・事業スナップショット → 理論 → ケース → 英文法　　所要時間：約{data['estimated_minutes']}分",
        config,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "ケースの会社・事業が何をしているかを理解してから、分析フレームを使う。", config)

    slide = prs.slides.add_slide(blank)
    snapshot_title = "会社スナップショット" if profile["profile_type"] == "company" else "事業スナップショット"
    snapshot_subtitle = "Company Context" if profile["profile_type"] == "company" else "Business Context"
    add_header(slide, f"{snapshot_title}：{profile['name']}", f"{snapshot_subtitle} | {profile['as_of']}", config)
    add_card(
        slide,
        0.45,
        0.98,
        4.02,
        1.22,
        f"何の会社・事業か\n{profile['overview']}",
        config,
        fill=PALE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        4.67,
        0.98,
        4.02,
        1.22,
        f"顧客と提供価値\n顧客：{profile['customers']}\n提供：{profile['offerings']}",
        config,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        8.89,
        0.98,
        3.99,
        1.22,
        f"ビジネスモデル\n{profile['business_model']}",
        config,
        fill=GREEN,
        line=LINE,
        font_size=config["output"]["note_pt"],
    )
    scale_text = "\n".join(profile["scale_indicators"])
    add_card(
        slide,
        0.45,
        2.50,
        4.02,
        1.45,
        f"規模・構造の目安\n{scale_text}",
        config,
        font_size=config["output"]["note_pt"],
    )
    value_chain = " → ".join(profile["value_chain"])
    add_card(
        slide,
        4.67,
        2.50,
        4.02,
        1.45,
        f"価値を生む流れ\n{value_chain}",
        config,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        8.89,
        2.50,
        3.99,
        1.45,
        f"今回見る理由\n{profile['why_this_case']}",
        config,
        fill=PALE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        0.45,
        4.28,
        12.43,
        1.38,
        f"このページを読んだ時点の仮説\n{profile['analysis_hypothesis']}",
        config,
        fill=GREEN,
        line=BLUE,
        bold=True,
    )
    source_text = f"出典：{profile['source_label']}　｜　確認日：{profile['verified_on']}"
    source = add_card(
        slide,
        0.45,
        5.95,
        12.43,
        0.55,
        source_text,
        config,
        font_size=config["output"]["note_pt"],
    )
    if profile.get("source_url"):
        for paragraph in source.text_frame.paragraphs:
            for run in paragraph.runs:
                run.hyperlink.address = profile["source_url"]
    add_footer(slide, "数字には時点を付ける。会社ケースと一般化した業界ケースを混同しない。", config)

    slide = prs.slides.add_slide(blank)
    add_header(slide, lesson["title"], "Theory", config)
    add_card(
        slide,
        0.45,
        0.98,
        12.43,
        1.00,
        f"基本概念\n{lesson['definition']}",
        config,
        fill=PALE,
        line=BLUE,
        bold=True,
    )
    for i, point in enumerate(lesson["learning_points"]):
        add_card(
            slide,
            0.45 + i * 4.18,
            2.25,
            3.92,
            1.38,
            f"ポイント{i + 1}\n{point}",
            config,
            fill=WHITE,
            line=LINE,
            font_size=config["output"]["note_pt"],
        )
    chain = lesson["cause_effect_chain"]
    x_positions = [0.55, 4.77, 8.99]
    for i, item in enumerate(chain):
        add_card(
            slide,
            x_positions[i],
            4.02,
            3.42,
            1.02,
            item,
            config,
            fill=PALE,
            line=BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
            font_size=config["output"]["note_pt"],
        )
        if i < 2:
            add_arrow(slide, x_positions[i] + 3.52, 4.32)
    add_card(
        slide,
        0.45,
        5.42,
        12.43,
        0.95,
        f"よくある混同\n{lesson['common_confusion']}",
        config,
        fill=WHITE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    add_footer(slide, "フレームワーク名より、顧客価値とコスト差が生まれる因果を説明する。", config)

    case = data["case_lesson"]
    slide = prs.slides.add_slide(blank)
    add_header(slide, case["title"], "Case Application | Evidence and hypothesis", config)
    add_card(
        slide,
        0.45,
        0.98,
        12.43,
        0.92,
        f"ケースの状況\n{case['context']}",
        config,
        fill=PALE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    obs_text = "\n".join(f"・{item}" for item in case["observations"])
    int_text = "\n".join(f"・{item}" for item in case["interpretation"])
    add_card(
        slide,
        0.45,
        2.18,
        5.95,
        2.35,
        f"観察できること\n{obs_text}",
        config,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        6.67,
        2.18,
        6.21,
        2.35,
        f"理論で読み解くと\n{int_text}",
        config,
        fill=WHITE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        0.45,
        4.85,
        12.43,
        1.25,
        f"このケースからの要点\n{case['takeaway']}",
        config,
        fill=GREEN,
        line=BLUE,
        bold=True,
        font_size=config["output"]["closing_pt"],
    )
    add_footer(slide, f"参照区分：{case['reference_label']}。観察事実と分析仮説を分ける。", config)

    english = data["english_lesson"]
    slide = prs.slides.add_slide(blank)
    add_header(slide, english["title"], f"Business English | {profile['name']}と接続", config)
    add_card(
        slide,
        0.45,
        0.98,
        12.43,
        0.75,
        f"今日の型：{english['grammar_pattern']}",
        config,
        fill=PALE,
        line=BLUE,
        bold=True,
        font_size=config["output"]["closing_pt"],
        align=PP_ALIGN.CENTER,
    )
    rules = "\n".join(f"{i + 1}. {item}" for i, item in enumerate(english["rule_points"]))
    add_card(slide, 0.45, 1.98, 4.05, 1.60, f"ルール\n{rules}", config, font_size=config["output"]["note_pt"])
    add_card(slide, 4.73, 1.98, 3.92, 1.60, f"誤り例\n{english['bad_example']}", config, font_size=config["output"]["note_pt"])
    add_card(
        slide,
        8.88,
        1.98,
        4.00,
        1.60,
        f"正解例\n{english['correct_example']}",
        config,
        fill=WHITE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
    )
    breakdown = " / ".join(english["example_breakdown"])
    add_card(
        slide,
        0.45,
        3.85,
        12.43,
        0.70,
        f"文の分解：{breakdown}",
        config,
        font_size=config["output"]["note_pt"],
    )
    vocab_text = " / ".join(f"{item['word']}（{item['meaning']}）" for item in english["vocabulary"])
    add_card(
        slide,
        0.45,
        4.78,
        12.43,
        0.78,
        f"今日の5語：{vocab_text}",
        config,
        fill=PALE,
        line=LINE,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        0.45,
        5.78,
        7.75,
        0.78,
        f"完成例文\n{english['worked_example_ja']}\n{english['worked_example_en']}",
        config,
        font_size=config["output"]["note_pt"],
    )
    add_card(
        slide,
        8.45,
        5.78,
        4.43,
        0.78,
        f"ミニ練習\n{english['practice_prompt']}",
        config,
        fill=WHITE,
        line=BLUE,
        font_size=config["output"]["note_pt"],
        bold=True,
    )
    add_footer(slide, "長く書くより、主語・動詞・因果を崩さない。", config)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    office = shutil.which("libreoffice") or shutil.which("soffice")
    if not office:
        raise RuntimeError("LibreOffice is required to export the lesson PDF")
    with tempfile.TemporaryDirectory(prefix="growth-pdf-") as temp_dir:
        result = subprocess.run(
            [office, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, str(pptx_path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=180,
        )
        generated = Path(temp_dir) / f"{pptx_path.stem}.pdf"
        if result.returncode != 0 or not generated.exists():
            raise RuntimeError(
                "LibreOffice PDF export failed: "
                f"returncode={result.returncode} stdout={result.stdout[-1000:]} stderr={result.stderr[-1000:]}"
            )
        pdf_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(generated, pdf_path)


def render_preview(pdf_path: Path, png_path: Path) -> None:
    document = fitz.open(pdf_path)
    if document.page_count < 1:
        raise RuntimeError("Generated PDF has no pages")
    page = document[0]
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
    png_path.parent.mkdir(parents=True, exist_ok=True)
    pixmap.save(png_path)
    document.close()


def write_outputs(data: dict[str, Any], config: dict[str, Any], output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "history").mkdir(parents=True, exist_ok=True)
    date_slug = data["date"]
    latest_json = output_dir / "latest.json"
    history_json = output_dir / "history" / f"{date_slug}.json"
    latest_pdf = output_dir / "latest.pdf"
    history_pdf = output_dir / "history" / f"{date_slug}_growth_lesson.pdf"
    preview = output_dir / "latest-preview.png"
    notification = output_dir / "latest_notification.md"

    text = json.dumps(data, ensure_ascii=False, indent=2) + "\n"
    latest_json.write_text(text, encoding="utf-8")
    history_json.write_text(text, encoding="utf-8")

    with tempfile.TemporaryDirectory(prefix="growth-lesson-") as temp_dir:
        pptx_path = Path(temp_dir) / "growth_lesson.pptx"
        pdf_path = Path(temp_dir) / "growth_lesson.pdf"
        render_pptx(data, config, pptx_path)
        convert_pptx_to_pdf(pptx_path, pdf_path)
        shutil.copy2(pdf_path, latest_pdf)
        shutil.copy2(pdf_path, history_pdf)

    render_preview(latest_pdf, preview)

    pdf_url = "https://raw.githubusercontent.com/purachinaevo7-cmyk/HOS/main/outputs/growth-menu/latest.pdf"
    preview_url = "https://raw.githubusercontent.com/purachinaevo7-cmyk/HOS/main/outputs/growth-menu/latest-preview.png"
    notification.write_text(
        f"## {data['notification']['headline']}\n\n"
        f"{data['notification']['summary']}\n\n"
        f"所要時間：{data['notification']['duration']}\n\n"
        f"[今日のPDF授業を開く]({pdf_url})\n\n"
        f"[1ページ目を画像で確認]({preview_url})\n\n"
        "PDFは5ページです：学習マップ／会社・事業スナップショット／理論／ケース／英文法。\n\n"
        "学習後は、このチャットに次の3点を返してください。\n"
        "1. 今日理解したこと\n"
        "2. ケースに対する自分の見解\n"
        "3. 5ページ目の英作文\n",
        encoding="utf-8",
    )
    return {
        "latest_json": latest_json,
        "history_json": history_json,
        "latest_pdf": latest_pdf,
        "history_pdf": history_pdf,
        "preview": preview,
        "notification": notification,
    }
